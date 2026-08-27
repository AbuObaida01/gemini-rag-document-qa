from pathlib import Path
from typing import Any

from pptx import Presentation

from app.services.extraction.base import BaseExtractor


class PPTXExtractor(BaseExtractor):
    """
    Extract text from PowerPoint PPTX presentations.

    Extracts:
    - Slide text
    - Text boxes
    - Titles
    - Tables
    - Bullet points
    """

    def extract(self, file_path: Path) -> dict[str, Any]:
        if not file_path.exists():
            raise FileNotFoundError(
                f"File not found: {file_path}"
            )

        if file_path.suffix.lower() != ".pptx":
            raise ValueError(
                "PPTXExtractor can only process PPTX files."
            )

        presentation = Presentation(file_path)

        slides = []
        all_text = []

        for slide_number, slide in enumerate(
            presentation.slides,
            start=1,
        ):
            slide_parts = []

            for shape in slide.shapes:
                # Normal text-containing shapes
                if hasattr(shape, "text"):
                    text = shape.text.strip()

                    if text:
                        slide_parts.append(text)

                # Tables
                if shape.has_table:
                    table = shape.table

                    rows = []

                    for row in table.rows:
                        cells = [
                            cell.text.strip()
                            for cell in row.cells
                        ]

                        rows.append(
                            " | ".join(cells)
                        )

                    table_text = "\n".join(rows).strip()

                    if table_text:
                        slide_parts.append(table_text)

            slide_text = "\n".join(slide_parts).strip()

            if not slide_text:
                continue

            slide_data = {
                "text": slide_text,
                "metadata": {
                    "slide": slide_number,
                },
            }

            slides.append(slide_data)
            all_text.append(slide_text)

        if not slides:
            raise ValueError(
                "No extractable text was found in the PPTX file."
            )

        return {
            "text": "\n\n".join(all_text),
            "slides": slides,
            "metadata": {
                "filename": file_path.name,
                "file_type": "pptx",
                "slide_count": len(presentation.slides),
            },
        }